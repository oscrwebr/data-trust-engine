from app.scanning.detectors import re
from app.scanning.regex_patterns import re
import pymupdf
from io import BytesIO
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


# Helper method for normalising text by deleting line breaks
def normalise_text(text: str):
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text).strip()


# Helper method to call appropriate extractor method based on provided file_extension (file type)
def extract_text_from_file(file_bytes: bytes, file_extension, str):
    extension = file_extension.lower().lstrip(".")

    # Match the file extension to call appropriate extractor method, otherwise raise error if provided file type is unsupported
    match extension:
        case "pdf":
            return extract_text_from_pdf(file_bytes=file_bytes)
        case "txt":
            return extract_text_from_txt(file_bytes=file_bytes)
        case "docx":
            return extract_text_from_docx(file_bytes=file_bytes)
        case "xlsx":
            return extract_text_from_xlsx(file_bytes=file_bytes)
        case "pptx":
            return extract_text_from_pptx(file_bytes=file_bytes)
        case _:
            raise ValueError(f"Unsupported file type for extraction: {file_extension}")


# Extract text from .pdf filebytes into dict (pdf files)
def extract_text_from_pdf(file_bytes: bytes) -> dict:
    file = pymupdf.open(stream=file_bytes, filetype="pdf")
    extracted_text = {}

    # Make page numbers 1 indexed, because user think in page 1, 2, 3 not 0, 1, 2
    for page_number in range(len(file)):
        page = file.load_page(page_number)
        text = normalise_text(page.get_text("text"))

        extracted_text[page_number + 1] = text

    file.close()
    return extracted_text


# Extract text from .txt filebytes (text files)
def extract_text_from_txt(file_bytes: bytes) -> dict:
    text = file_bytes.decode("utf-8", errors="ignore")

    # Only one page on a .txt document therefore return as just page number 1
    return {
        1: normalise_text(text)
    }


# Extract text from .docx filebytes (word documents)
def extract_text_from_docx(file_bytes: bytes) -> dict:
    # Turn file_bytes into Document object
    document = Document(BytesIO(file_bytes))

    text_parts = []

    # Extract paragraphs text
    for paragraph in document.paragraphs:
        if paragraph.text:
            text_parts.append(paragraph.text)
    
    # Extract tables text
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    text_parts.append(cell.text)

    # Join all text_parts into one variable, separated by spaces
    full_text = " ".join(text_parts)

    # .docx is an XML, not split up by 'pages' therefore return it all as one page
    return {
        1: normalise_text(full_text)
    }


# Extract text from .pptx filebytes (powerpoint documents)
def extract_text_from_pptx(file_bytes: bytes) -> dict:
    # Turn file bytes into Presentation object
    presentation = Presentation(BytesIO(file_bytes))
    extracted_text = {}

    # Iterate through the slides of the presentation document
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        # Append to slide_text table if there is a text box instance on the current slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)

        extracted_text[slide_index] = normalise_text(" ".join(slide_text))

    return extracted_text


# Extract text from .xlsx filebytes (excel spreadsheet)
def extract_text_from_xlsx(file_bytes: bytes) -> dict:
    # Turn file bytes into a Workbook object
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)
    extracted_text = {}

    sheet_number = 1
    # Iterate through sheets of spreadsheet
    for sheet in workbook.worksheets:
        cell_values = []

        # Append data to cell_values table only if cell value is not empty
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    cell_values.append(str(cell))

        # Append this sheet's text into dictionary with current sheet number as key
        extracted_text[sheet_number] = normalise_text(" ".join(cell_values))
        sheet_number += 1

    return extracted_text